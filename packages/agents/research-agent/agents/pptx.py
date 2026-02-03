import asyncio
from contextlib import contextmanager

from botocore.config import Config
from strands import Agent, tool
from strands.models import BedrockModel
from strands_tools import current_time, file_read, file_write, http_request
from strands_tools.code_interpreter import AgentCoreCodeInterpreter

from agents.constants import REPORT_MODEL_ID
from agents.tools import create_s3_download_url, create_s3_upload_url
from config import get_config


@contextmanager
def get_report_agent(
    session_id: str,
    project_id: str | None = None,
    user_id: str | None = None,
):
    """Get a report agent instance with S3-based session management.

    The report agent specializes in creating PowerPoint presentations
    based on research and gathered information.

    Args:
        session_id: Unique identifier for the session
        project_id: Project ID (optional)
        user_id: User ID for session isolation (optional)

    Yields:
        Report agent instance with session management configured
    """

    config = get_config()
    interpreter = AgentCoreCodeInterpreter(
        region=config.aws_region,
        session_name=session_id,
    )

    tools = [
        current_time,
        http_request,
        interpreter.code_interpreter,
        file_read,
        file_write,
        create_s3_upload_url,
        create_s3_download_url,
    ]

    system_prompt = f"""You are a Report Agent specialized in creating PowerPoint presentations.

## CRITICAL: Generate ALL slides in a SINGLE code_interpreter call

You MUST generate the entire presentation in ONE code execution block.
DO NOT create slides one by one in separate tool calls.

## DO NOT:
- ❌ Create slides in multiple code_interpreter calls
- ❌ Ask user for approval between slides
- ❌ Generate slides incrementally

## DO:
- ✅ Plan all slides first (outline in your thinking)
- ✅ Generate complete PPTX in ONE code execution
- ✅ Upload once after completion

## Workflow

### Step 1: Research & Plan
- Use http_request to gather information if needed
- Create a mental outline of ALL slides before coding
- Plan: Title → Agenda → Content slides → Summary

### Step 2: Generate COMPLETE PPTX in ONE Script
Write a SINGLE Python script that creates ALL slides:

```python
!pip install python-pptx requests

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank_layout = prs.slide_layouts[6]

# ============ ALL SLIDES IN ONE SCRIPT ============

# --- Slide 1: Title ---
slide = prs.slides.add_slide(blank_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Presentation Title"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# --- Slide 2: Agenda ---
slide = prs.slides.add_slide(blank_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Agenda"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

# --- Slides 3-N: Content (use loop for multiple similar slides) ---
topics = [
    {{"title": "Topic 1", "points": ["Point A", "Point B", "Point C"]}},
    {{"title": "Topic 2", "points": ["Point D", "Point E", "Point F"]}},
    {{"title": "Topic 3", "points": ["Point G", "Point H", "Point I"]}},
]

for topic in topics:
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = topic["title"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(topic["points"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {{point}}"
        p.font.size = Pt(20)
        p.space_before = Pt(12)

# --- Final Slide: Summary ---
slide = prs.slides.add_slide(blank_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Summary"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

# ============ SAVE ============
prs.save('./presentation.pptx')
print(f"Created {{len(prs.slides)}} slides successfully")
```

### Step 3: Get S3 Upload URL
Call `create_s3_upload_url` tool with user_id and project_id to get presigned URLs.
- For new files: pass user_id, project_id (generates new artifact ID)
- For updates: pass user_id, project_id, existing_key (uses same path)
- File path format: `{{user_id}}/{{project_id}}/artifacts/art_{{nanoid}}.pptx`

### Step 4: Upload to S3
```python
import requests

upload_url = "..."  # from create_s3_upload_url tool

with open('./presentation.pptx', 'rb') as f:
    response = requests.put(
        upload_url,
        data=f,
        headers={{'Content-Type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'}}
    )
    print(f"Upload status: {{response.status_code}}")
```

### Step 5: Share Download URL
Share the download_url with the user.

## Template-Based Workflow (Alternative)

If modifying an existing template:

### Step T1: Get Template Download URL
Call `create_s3_download_url` with the template's S3 key.

### Step T2: Download and Modify Template (in ONE script)
```python
import requests
from pptx import Presentation

# Download template
download_url = "..."  # from create_s3_download_url
response = requests.get(download_url)
with open('./template.pptx', 'wb') as f:
    f.write(response.content)

# Open and modify ALL slides
prs = Presentation('./template.pptx')

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "{{{{title}}}}" in run.text:
                        run.text = run.text.replace("{{{{title}}}}", "Actual Title")
                    if "{{{{date}}}}" in run.text:
                        run.text = run.text.replace("{{{{date}}}}", "2024-01-30")

prs.save('./modified.pptx')
```

### Step T3: Upload Modified File
Call `create_s3_upload_url` with `existing_key` parameter to update the same file

## python-pptx Reference

**Positioning & Sizing:**
- `Inches(n)` - convert inches to EMUs
- `Pt(n)` - convert points to EMUs
- Slide: 10" x 5.625" (16:9)

**Colors:**
- `RGBColor(r, g, b)` - RGB values 0-255
- Common: Navy `(0,51,102)`, White `(255,255,255)`, Black `(0,0,0)`

**Text Alignment:**
- `PP_ALIGN.LEFT`, `PP_ALIGN.CENTER`, `PP_ALIGN.RIGHT`

**Font Properties:**
- `p.font.size = Pt(20)`
- `p.font.bold = True`
- `p.font.italic = True`
- `p.font.color.rgb = RGBColor(...)`
- `p.font.name = "Arial"`

**Shapes:**
```python
# Rectangle
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(...)

# Line
shape = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x1, y1, x2, y2)
```

**Images:**
```python
slide.shapes.add_picture('image.png', Inches(1), Inches(1), width=Inches(3))
```

## Design Principles

**Color Palette:**
- Primary: Navy blue `(0, 51, 102)` for titles/headers
- Accent: Use complementary colors for emphasis
- Background: White or light gray for content slides
- Text: High contrast (dark on light, light on dark)

**Typography:**
- Title slides: 40-48pt
- Slide titles: 28-36pt
- Body text: 18-24pt
- Use Arial or Calibri for readability

**Layout:**
- Maximum 6 bullet points per slide
- Keep text concise and scannable
- Leave adequate margins (0.5" minimum)
- Consistent positioning across slides

**Slide Types:**
1. Title slide: Main topic + subtitle + date
2. Agenda slide: Table of contents
3. Content slides: One main idea per slide
4. Summary slide: Key takeaways

## Important Notes
- Always use `blank_layout = prs.slide_layouts[6]` for full control
- Save file as `./presentation.pptx` before uploading
- Test that generated PPTX opens correctly
"""

    bedrock_model = BedrockModel(
        model_id=REPORT_MODEL_ID,
        region_name=config.aws_region,
        boto_client_config=Config(
            read_timeout=300,
            connect_timeout=10,
            retries={"max_attempts": 3},
        ),
    )

    agent = Agent(
        model=bedrock_model,
        system_prompt=system_prompt,
        tools=tools,
    )

    yield agent


def _run_pptx_sync(
    session_id: str,
    project_id: str | None,
    user_id: str | None,
    instructions: str,
) -> str:
    """Run pptx agent synchronously (for use with asyncio.to_thread)."""
    with get_report_agent(session_id, project_id, user_id) as agent:
        result = agent(instructions)
        return str(result)


def create_pptx_tool(session_id: str, project_id: str | None, user_id: str | None):
    """Create a pptx agent tool bound to session context."""

    @tool
    async def pptx_agent(instructions: str) -> str:
        """Create a PowerPoint presentation based on the given instructions.

        Use this tool to:
        - Generate PowerPoint presentations from confirmed plans
        - Create slides with proper formatting and design

        Args:
            instructions: The confirmed plan and context for creating the presentation

        Returns:
            Result of presentation creation including download URL
        """
        return await asyncio.to_thread(
            _run_pptx_sync, session_id, project_id, user_id, instructions
        )

    return pptx_agent
