"""Extract tables from PPTX using python-pptx."""

import io

from pptx import Presentation

from .artifact import get_artifact_content


def table_to_markdown(table) -> str:
    """Convert a PPTX table to markdown format."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    lines = []

    # Header row
    header = rows[0]
    lines.append("| " + " | ".join(header) + " |")

    # Separator
    lines.append("| " + " | ".join("---" for _ in header) + " |")

    # Data rows
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def extract_tables(event: dict) -> dict:
    """Extract tables from PPTX.

    Args:
        event: {
            "artifact_id": str,
            "format": "json" | "markdown"  # default: json
        }

    Returns:
        {
            "tables": list[{
                "slide_number": int,
                "table_index": int,
                "data": list[list[str]] | str  # depends on format
            }]
        }
    """
    artifact_id = event["artifact_id"]
    output_format = event.get("format", "json")

    content, metadata = get_artifact_content(artifact_id)

    if not metadata.content_type.endswith(("presentationml.presentation", "/pptx")):
        raise ValueError(f"Artifact is not a PPTX: {metadata.content_type}")

    pptx_file = io.BytesIO(content)
    prs = Presentation(pptx_file)

    tables_result = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        table_index = 0

        for shape in slide.shapes:
            if not shape.has_table:
                continue

            table = shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)

            if not rows:
                continue

            data = table_to_markdown(table) if output_format == "markdown" else rows

            tables_result.append({
                "slide_number": slide_num,
                "table_index": table_index,
                "data": data,
            })
            table_index += 1

    return {"tables": tables_result}
