"""Extract text from PPTX using python-pptx."""

import io

from pptx import Presentation

from .artifact import get_artifact_content


def extract_text(event: dict) -> dict:
    """Extract text from PPTX.

    Args:
        event: {
            "artifact_id": str
        }

    Returns:
        {
            "text": str,
            "slides": list[{"slide_number": int, "text": str, "notes": str}]
        }
    """
    artifact_id = event["artifact_id"]

    content, metadata = get_artifact_content(artifact_id)

    if not metadata.content_type.endswith(("presentationml.presentation", "/pptx")):
        raise ValueError(f"Artifact is not a PPTX: {metadata.content_type}")

    pptx_file = io.BytesIO(content)
    prs = Presentation(pptx_file)

    slides_result = []
    all_text_parts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []

        # Extract text from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text_parts.append(text)

            # Extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            slide_text_parts.append(text)

        slide_text = "\n".join(slide_text_parts)

        # Extract notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()

        slides_result.append({
            "slide_number": slide_num,
            "text": slide_text,
            "notes": notes_text,
        })

        if slide_text:
            all_text_parts.append(f"[Slide {slide_num}]\n{slide_text}")
        if notes_text:
            all_text_parts.append(f"[Notes {slide_num}]\n{notes_text}")

    return {
        "text": "\n\n".join(all_text_parts),
        "slides": slides_result,
    }
